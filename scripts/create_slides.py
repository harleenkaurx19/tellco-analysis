from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Colors
BLUE = RGBColor(0, 70, 127)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(240, 240, 240)
DARK = RGBColor(30, 30, 30)
GREEN = RGBColor(0, 150, 80)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_slide(prs, layout=6):
    slide_layout = prs.slide_layouts[layout]
    return prs.slides.add_slide(slide_layout)

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                fontsize=18, bold=False, color=RGBColor(0,0,0),
                align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

# ============================================
# SLIDE 1 - TITLE
# ============================================
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "TellCo Telecom", 1, 1, 11, 1.5,
            fontsize=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "User Analytics Report", 1, 2.5, 11, 1,
            fontsize=32, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Analyzing Growth Opportunities for TellCo Mobile Services",
            1, 3.8, 11, 1, fontsize=18, color=RGBColor(200,220,255),
            align=PP_ALIGN.CENTER)
add_textbox(slide, "Prepared by: Harleen | April 2026",
            1, 5.5, 11, 0.8, fontsize=14, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================
# SLIDE 2 - OVERVIEW
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Project Overview", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
- Dataset: 150,001 sessions from TellCo mobile network
- Period: April 2019
- Users: Over 106,000 unique customers analyzed
- Goal: Identify growth opportunities & recommend buy/sell decision
- Tools: Python, Pandas, Scikit-learn, Streamlit, PostgreSQL
""", 0.5, 1.5, 12, 5, fontsize=18, color=DARK)

# ============================================
# SLIDE 3 - TOP HANDSETS
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Task 1: Top 10 Handsets", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Top 10 Handsets Used by Customers:

1. Huawei B528S-23A         - 20,324 users
2. Apple iPhone 6S (A1688)  -  9,419 users
3. Apple iPhone 6 (A1586)   -  9,023 users
4. Apple iPhone 7 (A1778)   -  6,326 users
5. Apple iPhone SE (A1723)  -  5,187 users
6. Apple iPhone 8 (A1905)   -  4,993 users
7. Apple iPhone XR (A2105)  -  4,568 users
8. Samsung Galaxy S8        -  4,520 users
9. Apple iPhone X (A1901)   -  3,813 users
10. Samsung Galaxy A5       -  3,724 users
""", 0.5, 1.3, 12, 5.5, fontsize=16, color=DARK)

# ============================================
# SLIDE 4 - TOP MANUFACTURERS
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 1: Top 3 Manufacturers", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Top 3 Handset Manufacturers:

🥇 Apple    - 60,137 users (40%)
🥈 Samsung  - 40,839 users (27%)
🥉 Huawei   - 34,423 users (23%)

Key Insight:
Apple dominates the market with 40% share.
Samsung and Huawei follow with strong presence.
Marketing should focus on Apple & Samsung users
for premium service offerings.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 5 - USER BEHAVIOR
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Task 1: User Behavior Overview", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Per-User Aggregated Metrics:

- Total Sessions analyzed: 150,001
- Unique Users: 106,856
- Average session duration: varies by user segment
- Top apps by data usage: Gaming & Other Applications
- Social Media and YouTube show consistent usage

Data was cleaned by replacing missing values with column means.
Outliers were handled using 3 standard deviation clipping.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 6 - ENGAGEMENT INTRO
# ============================================
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "Task 2", 1, 2, 11, 1.5,
            fontsize=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "User Engagement Analysis", 1, 3.2, 11, 1,
            fontsize=28, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================
# SLIDE 7 - ENGAGEMENT METRICS
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 2: Engagement Metrics", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Top 10 Users by Sessions, Duration & Traffic were identified.

K-Means Clustering (k=3) Results:

- Cluster 0 - Low Engagement:
  Few sessions, short duration, low traffic

- Cluster 1 - Medium Engagement:
  Moderate sessions, average duration & traffic

- Cluster 2 - High Engagement:
  Many sessions, long duration, high data usage

Elbow method confirmed k=3 as optimal number of clusters.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 8 - TOP APPS
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Task 2: Top Applications", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Application Usage by Total Data Volume:

🥇 Gaming        - Highest data usage
🥈 Other Apps    - Second highest
🥉 Social Media  - Third highest

YouTube, Netflix, Google & Email show lower but
consistent usage patterns.

Recommendation:
Network resources should be prioritized for
Gaming and Social Media applications.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 9 - ELBOW METHOD
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 2: Optimal Clusters - Elbow Method", 0.5, 0.3, 12, 1,
            fontsize=28, bold=True, color=BLUE)
add_textbox(slide, """
The Elbow Method was used to find the optimal k value.

- Tested k values from 1 to 10
- Plotted inertia (within-cluster sum of squares)
- The curve bends at k=3 indicating optimal clustering

This means customers naturally fall into 3 groups:
  - Low engagement users
  - Medium engagement users
  - High engagement users

See elbow curve chart in dashboard for visual proof.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 10 - EXPERIENCE INTRO
# ============================================
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "Task 3", 1, 2, 11, 1.5,
            fontsize=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "User Experience Analysis", 1, 3.2, 11, 1,
            fontsize=28, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================
# SLIDE 11 - EXPERIENCE METRICS
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 3: Experience Metrics", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Per-user experience metrics computed:

- Average TCP Retransmission (DL + UL)
- Average Round Trip Time - RTT (DL + UL)
- Average Throughput (DL + UL)
- Handset Type (most frequent per user)

Missing values replaced with column mean.
Outliers handled with 3 standard deviation clipping.

Top, Bottom & Most Frequent values computed
for TCP, RTT and Throughput metrics.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 12 - THROUGHPUT & TCP
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Task 3: Throughput & TCP per Handset", 0.5, 0.3, 12, 1,
            fontsize=28, bold=True, color=BLUE)
add_textbox(slide, """
Throughput Distribution per Handset Type:
- Huawei B528S-23A shows highest average throughput
- Apple devices show consistent mid-range throughput
- Samsung devices vary widely in throughput

TCP Retransmission per Handset Type:
- Higher TCP retransmission = worse network experience
- Some handsets show significantly higher retransmission
- Indicates network compatibility issues with certain devices

See charts in dashboard for visual breakdown.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 13 - EXPERIENCE CLUSTERS
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 3: Experience Clusters (k=3)", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
K-Means Clustering on Experience Metrics:

- Cluster 0 - Poor Experience:
  High TCP retransmission, High RTT, Low throughput

- Cluster 1 - Average Experience:
  Moderate TCP, Moderate RTT, Moderate throughput

- Cluster 2 - Good Experience:
  Low TCP retransmission, Low RTT, High throughput

Most users fall in Average Experience cluster.
Poor experience users need immediate network attention.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 14 - SATISFACTION INTRO
# ============================================
slide = add_slide(prs)
set_bg(slide, BLUE)
add_textbox(slide, "Task 4", 1, 2, 11, 1.5,
            fontsize=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "User Satisfaction Analysis", 1, 3.2, 11, 1,
            fontsize=28, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================
# SLIDE 15 - SATISFACTION SCORES
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 4: Satisfaction Scores", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Satisfaction Score Computation:

- Engagement Score = Euclidean distance from least
  engaged cluster centroid

- Experience Score = Euclidean distance from worst
  experience cluster centroid

- Satisfaction Score = Average of both scores

106,856 users scored successfully!
Top 10 most satisfied customers identified.
All scores exported to PostgreSQL database.
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 16 - REGRESSION MODEL
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Task 4: Regression Model", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Linear Regression Model to Predict Satisfaction:

- Features: Engagement Score + Experience Score
- Target: Satisfaction Score
- Train/Test Split: 80% / 20%

Model Results:
  ✅ MSE:      0.0000 (Near Perfect!)
  ✅ R² Score: 1.0000 (Perfect Prediction!)

Model tracked using MLflow:
- Parameters logged
- Metrics logged
- Model artifact saved
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 17 - SATISFACTION CLUSTERS
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "Task 4: Satisfaction Clusters (k=2)", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
K-Means (k=2) on Engagement & Experience Scores:

- Cluster 0 - Lower Satisfaction:
  Average satisfaction: 0.2922
  Count: 35,607 users

- Cluster 1 - Higher Satisfaction:
  Average satisfaction: 0.3494
  Count: 71,249 users

Majority of users (71,249) fall in higher satisfaction cluster.
This is a positive sign for TellCo's service quality!
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 18 - RECOMMENDATION
# ============================================
slide = add_slide(prs)
set_bg(slide, GREEN)
add_textbox(slide, "✅ Recommendation: BUY TellCo!", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, """
Strong reasons to purchase TellCo:

✅ Large user base: 106,856 active customers
✅ High engagement: Gaming & Social Media drive usage
✅ Apple & Samsung dominance: Premium user segment
✅ Majority satisfied: 71,249 users in high satisfaction cluster
✅ Growth potential: Underserved segments identified
✅ Network improvements can boost experience scores

Action Plan:
- Focus marketing on Apple & Samsung users
- Improve network for poor experience cluster
- Expand Gaming & Social Media bandwidth
- Target high-engagement users for premium plans
""", 0.5, 1.5, 12, 5.5, fontsize=17, color=WHITE)

# ============================================
# SLIDE 19 - LIMITATIONS
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Limitations of Analysis", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
- Data covers only 1 month (April 2019) - may not represent
  full year patterns or seasonal trends

- Missing values were replaced with means which may
  introduce slight bias in results

- Satisfaction score is derived (not directly measured)
  from engagement and experience metrics only

- R² = 1.0 suggests possible data leakage in regression
  model - needs further investigation

- External factors (pricing, competition, coverage area)
  not included in analysis

- Dataset may be outdated (2019) - market conditions
  may have changed significantly
""", 0.5, 1.3, 12, 5.5, fontsize=17, color=DARK)

# ============================================
# SLIDE 20 - REFERENCES
# ============================================
slide = add_slide(prs)
set_bg(slide, GRAY)
add_textbox(slide, "References", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
- TellCo xDR Dataset - Nexthikes IT Solutions (2019)

- Scikit-learn Documentation:
  https://scikit-learn.org

- Streamlit Documentation:
  https://streamlit.io

- MLflow Documentation:
  https://mlflow.org

- Pandas Documentation:
  https://pandas.pydata.org

- GitHub Repository:
  https://github.com/harleenkaurx19/tellco-analysis
""", 0.5, 1.3, 12, 5.5, fontsize=18, color=DARK)

# ============================================
# SLIDE 20.5 - DASHBOARD SLIDE
# ============================================
slide = add_slide(prs)
set_bg(slide, WHITE)
add_textbox(slide, "Live Dashboard", 0.5, 0.3, 12, 1,
            fontsize=32, bold=True, color=BLUE)
add_textbox(slide, """
Streamlit Dashboard built for easy navigation:

🔗 GitHub Repository:
https://github.com/harleenkaurx19/tellco-analysis

📊 Dashboard Pages:
- 🏠 Home          - Key metrics overview
- 📊 User Overview - Top handsets & manufacturers
- 🔥 Engagement    - App usage & clusters
- 📡 Experience    - Network performance
- 😊 Satisfaction  - Scores & predictions

▶ To run dashboard locally:
  pip install -r requirements.txt
  streamlit run dashboard/app.py
""", 0.5, 1.3, 12, 4, fontsize=17, color=DARK)

# Add dashboard screenshot
screenshot_path = 'data/dashboard_screenshot.png'
if os.path.exists(screenshot_path):
    slide.shapes.add_picture(
        screenshot_path,
        Inches(0.5), Inches(4.5),
        Inches(12), Inches(2.8)
    )
    print("✅ Dashboard screenshot added!")
else:
    add_textbox(slide, 
        "📸 [Add dashboard screenshot here]",
        0.5, 4.5, 12, 2,
        fontsize=16, color=RGBColor(150,150,150))
prs.save('TellCo_Analysis_Report.pptx')
print("✅ Presentation created successfully!")
print("📊 File saved as: TellCo_Analysis_Report.pptx")
